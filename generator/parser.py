"""Parse PowerPoint files into structured slide data."""

from __future__ import annotations

import base64
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _extract_bullets(shape) -> list[str]:
    if not shape.has_text_frame:
        return []
    bullets: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = _clean_text(paragraph.text)
        if text:
            bullets.append(text)
    return bullets


def _extract_shape_content(shape) -> dict[str, Any] | None:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        return {
            "type": "image",
            "filename": image.filename or "image.png",
            "content_type": image.content_type,
            "data_base64": base64.b64encode(image.blob).decode("ascii"),
        }

    if shape.has_text_frame:
        bullets = _extract_bullets(shape)
        if bullets:
            return {"type": "text", "bullets": bullets}

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        children: list[dict[str, Any]] = []
        for child in shape.shapes:
            content = _extract_shape_content(child)
            if content:
                children.append(content)
        if children:
            return {"type": "group", "children": children}

    if shape.has_table:
        rows: list[list[str]] = []
        for row in shape.table.rows:
            rows.append([_clean_text(cell.text) for cell in row.cells])
        return {"type": "table", "rows": rows}

    return None


def _guess_title(shapes_content: list[dict[str, Any]]) -> str:
    for item in shapes_content:
        if item.get("type") == "text" and item.get("bullets"):
            return item["bullets"][0]
    return "Untitled Slide"


def parse_pptx(pptx_path: Path) -> dict[str, Any]:
    """Extract slides, notes, and assets from a PowerPoint file."""
    presentation = Presentation(str(pptx_path))
    slides: list[dict[str, Any]] = []

    for index, slide in enumerate(presentation.slides, start=1):
        shapes_content: list[dict[str, Any]] = []
        for shape in slide.shapes:
            content = _extract_shape_content(shape)
            if content:
                shapes_content.append(content)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = _clean_text(slide.notes_slide.notes_text_frame.text)

        slides.append(
            {
                "index": index,
                "title": _guess_title(shapes_content),
                "shapes": shapes_content,
                "notes": notes,
            }
        )

    core_props = presentation.core_properties
    return {
        "title": core_props.title or pptx_path.stem,
        "author": core_props.author or "",
        "slide_count": len(slides),
        "slides": slides,
    }
